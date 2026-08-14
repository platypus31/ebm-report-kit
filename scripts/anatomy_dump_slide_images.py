#!/usr/bin/env python3
"""範本解剖工具②：抽指定 slide 的所有圖片（挑固定資產候選用；挑選候選圖用（僅供人工參考））。
用法：python3 scripts/anatomy_dump_slide_images.py <pptx> <slide#,slide#…> <outdir>
"""
import sys, os
from pptx import Presentation

def walk(shapes):
    for sh in shapes:
        if sh.shape_type == 6:
            yield from walk(sh.shapes)
        elif sh.shape_type == 13:  # picture
            yield sh

path, nums, outdir = sys.argv[1], [int(x) for x in sys.argv[2].split(",")], sys.argv[3]
os.makedirs(outdir, exist_ok=True)
prs = Presentation(path)
tag = os.path.splitext(os.path.basename(path))[0][:12].replace(" ", "_")
for i, slide in enumerate(prs.slides, 1):
    if i not in nums:
        continue
    for j, pic in enumerate(walk(slide.shapes), 1):
        try:
            img = pic.image
        except Exception as e:  # linked（非內嵌）圖片無 blob，跳過不中斷批次
            print(f"⚠️ s{i} 第{j}張圖無法讀取（linked image?）：{e}", file=sys.stderr)
            continue
        w = pic.width or 0
        fn = f"{outdir}/{tag}-s{i:02d}-p{j}-{w//9525}px.{img.ext}"
        with open(fn, "wb") as f:
            f.write(img.blob)
        print(fn, len(img.blob) // 1024, "KB")
