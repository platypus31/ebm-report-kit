#!/usr/bin/env python3
"""範本解剖工具①：dump 一份 pptx 的逐頁文字大綱（標題/內文/表格/備忘稿/圖片佔位）。
新增院方範本要併入 data/report-spec.md 時先跑這支，逐頁對照 §1 不變量。
用法：python3 scripts/anatomy_dump_pptx.py <pptx> > out.txt
"""
import sys
from pptx import Presentation

def shape_text(shape, depth=0):
    out = []
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            out += shape_text(s, depth)
        return out
    if shape.has_table:
        tbl = shape.table
        rows = []
        for r in tbl.rows:
            cells = [c.text.strip().replace("\n", "⏎") for c in r.cells]
            rows.append(" | ".join(cells))
        out.append("  [表格]\n    " + "\n    ".join(rows))
        return out
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                out.append("  " * (para.level + 1) + t)
    if shape.shape_type == 13:  # picture
        out.append("  [圖片]")
    if getattr(shape, "has_chart", False):
        out.append("  [圖表 chart]")
    return out

path = sys.argv[1]
prs = Presentation(path)
print(f"=== {path}  ({len(prs.slides)} slides, {prs.slide_width}x{prs.slide_height}) ===")
for i, slide in enumerate(prs.slides, 1):
    layout = slide.slide_layout.name
    print(f"\n--- Slide {i} [{layout}] ---")
    for shape in slide.shapes:
        for line in shape_text(shape):
            print(line)
    if slide.has_notes_slide:
        nt = slide.notes_slide.notes_text_frame.text.strip()
        if nt:
            print(f"  [備忘稿] {nt[:300]}")
